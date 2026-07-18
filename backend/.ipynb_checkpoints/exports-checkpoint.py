"""Export dataset + insights to PDF, PPTX, and XLSX."""
from __future__ import annotations

import io
from typing import Any

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt
from reportlab.lib import colors
from reportlab.lib.pagesizes import letter
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.units import inch
from reportlab.platypus import (
    Image,
    Paragraph,
    SimpleDocTemplate,
    Spacer,
    Table,
    TableStyle,
)

PASTEL_BARS = ["#A7D8B7", "#F4B7AE", "#C7B7EF", "#F5DC85", "#9FC5E8"]


def _chart_to_png(chart: dict[str, Any]) -> bytes:
    """Render one chart spec to a PNG bytes buffer using matplotlib."""
    fig, ax = plt.subplots(figsize=(6.5, 3.2), dpi=140)
    data = chart.get("data", [])
    xs = [d["x"] for d in data]
    ys = [d["y"] or 0 for d in data]
    ctype = chart.get("type")
    if ctype == "line":
        ax.plot(xs, ys, color="#3F3F46", linewidth=2.4, marker="o", markersize=4)
        ax.fill_between(range(len(xs)), ys, alpha=0.08, color="#3F3F46")
    elif ctype == "bar":
        ax.bar(xs, ys, color=PASTEL_BARS[0])
    elif ctype == "histogram":
        ax.bar(xs, ys, color=PASTEL_BARS[2])
    elif ctype == "pie":
        colors_cycle = [PASTEL_BARS[i % len(PASTEL_BARS)] for i in range(len(xs))]
        ax.pie(
            ys,
            labels=xs,
            colors=colors_cycle,
            autopct=lambda p: f"{p:.0f}%" if p >= 5 else "",
            textprops={"fontsize": 8, "color": "#3F3F46"},
            wedgeprops={"linewidth": 1, "edgecolor": "white"},
        )
        ax.set_aspect("equal")
    if ctype != "pie":
        ax.set_xlabel(chart.get("x_label", ""), fontsize=9, color="#71717A")
        ax.set_ylabel(chart.get("y_label", ""), fontsize=9, color="#71717A")
        ax.tick_params(axis="x", labelrotation=30, labelsize=8)
        ax.tick_params(axis="y", labelsize=8)
        for spine in ("top", "right"):
            ax.spines[spine].set_visible(False)
        for spine in ("bottom", "left"):
            ax.spines[spine].set_color("#E4E4E7")
        ax.grid(axis="y", linestyle="--", alpha=0.35, color="#E4E4E7")
    ax.set_title(chart.get("title", ""), fontsize=12, fontweight="600", color="#18181B", loc="left")
    fig.tight_layout()
    buf = io.BytesIO()
    fig.savefig(buf, format="png", bbox_inches="tight")
    plt.close(fig)
    buf.seek(0)
    return buf.getvalue()


def build_pdf(name: str, summary: dict, insights: list[str], charts: list[dict], preview: list[dict]) -> bytes:
    buf = io.BytesIO()
    doc = SimpleDocTemplate(buf, pagesize=letter, leftMargin=54, rightMargin=54, topMargin=54, bottomMargin=48)
    styles = getSampleStyleSheet()
    styles.add(ParagraphStyle(name="H1c", parent=styles["Heading1"], textColor=colors.HexColor("#18181B"), fontSize=20))
    styles.add(ParagraphStyle(name="H2c", parent=styles["Heading2"], textColor=colors.HexColor("#18181B"), fontSize=13, spaceBefore=14))
    body = ParagraphStyle(name="Bodyc", parent=styles["BodyText"], textColor=colors.HexColor("#3F3F46"), fontSize=10, leading=14)

    story = []
    story.append(Paragraph(f"Data Report - {name}", styles["H1c"]))
    story.append(Spacer(1, 8))
    story.append(Paragraph(
        f"{summary.get('rows', 0):,} rows across {summary.get('columns', 0)} columns.",
        body,
    ))
    if summary.get("date_range"):
        dr = summary["date_range"]
        story.append(Paragraph(f"Time range ({dr['column']}): {dr['start'][:10]} to {dr['end'][:10]}", body))

    story.append(Paragraph("Key Insights", styles["H2c"]))
    for ins in insights or ["No insights available."]:
        story.append(Paragraph(f"- {ins}", body))

    for ch in charts:
        story.append(Paragraph(ch.get("title", "Chart"), styles["H2c"]))
        img_bytes = _chart_to_png(ch)
        img = Image(io.BytesIO(img_bytes), width=6.5 * inch, height=3.0 * inch)
        story.append(img)

    if preview:
        story.append(Paragraph("Data Preview (first rows)", styles["H2c"]))
        cols = list(preview[0].keys())[:6]
        table_data = [cols] + [[str(row.get(c, ""))[:24] for c in cols] for row in preview[:8]]
        tbl = Table(table_data, hAlign="LEFT")
        tbl.setStyle(TableStyle([
            ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor("#F4F4F5")),
            ("TEXTCOLOR", (0, 0), (-1, 0), colors.HexColor("#18181B")),
            ("GRID", (0, 0), (-1, -1), 0.4, colors.HexColor("#E4E4E7")),
            ("FONTSIZE", (0, 0), (-1, -1), 8),
            ("FONTNAME", (0, 0), (-1, 0), "Helvetica-Bold"),
            ("VALIGN", (0, 0), (-1, -1), "MIDDLE"),
        ]))
        story.append(tbl)

    doc.build(story)
    return buf.getvalue()


def build_pptx(name: str, summary: dict, insights: list[str], charts: list[dict]) -> bytes:
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    blank = prs.slide_layouts[6]

    # Title slide
    slide = prs.slides.add_slide(blank)
    tb = slide.shapes.add_textbox(Inches(0.6), Inches(2.6), Inches(12), Inches(2))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.text = "Data Report"
    p.font.size = Pt(44)
    p.font.bold = True
    p2 = tf.add_paragraph()
    p2.text = name
    p2.font.size = Pt(22)
    p3 = tf.add_paragraph()
    p3.text = f"{summary.get('rows', 0):,} rows \u2022 {summary.get('columns', 0)} columns"
    p3.font.size = Pt(16)

    # Insights slide
    slide = prs.slides.add_slide(blank)
    tb = slide.shapes.add_textbox(Inches(0.6), Inches(0.5), Inches(12), Inches(1))
    tb.text_frame.text = "Key Insights"
    tb.text_frame.paragraphs[0].font.size = Pt(30)
    tb.text_frame.paragraphs[0].font.bold = True
    body_box = slide.shapes.add_textbox(Inches(0.6), Inches(1.6), Inches(12), Inches(5.5))
    body_tf = body_box.text_frame
    body_tf.word_wrap = True
    for i, ins in enumerate(insights or ["No insights available."]):
        para = body_tf.paragraphs[0] if i == 0 else body_tf.add_paragraph()
        para.text = f"- {ins}"
        para.font.size = Pt(18)

    # Chart slides
    for ch in charts:
        slide = prs.slides.add_slide(blank)
        title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(12), Inches(0.7))
        title_box.text_frame.text = ch.get("title", "Chart")
        title_box.text_frame.paragraphs[0].font.size = Pt(24)
        title_box.text_frame.paragraphs[0].font.bold = True
        img_bytes = _chart_to_png(ch)
        img_stream = io.BytesIO(img_bytes)
        slide.shapes.add_picture(img_stream, Inches(0.8), Inches(1.3), width=Inches(11.5))

    out = io.BytesIO()
    prs.save(out)
    return out.getvalue()


def build_xlsx(df: pd.DataFrame, summary: dict, insights: list[str]) -> bytes:
    out = io.BytesIO()
    with pd.ExcelWriter(out, engine="xlsxwriter") as writer:
        df.to_excel(writer, sheet_name="Data", index=False)
        pd.DataFrame({"Insight": insights or ["No insights available."]}).to_excel(
            writer, sheet_name="Insights", index=False
        )
        summary_rows = [
            ["Rows", summary.get("rows", 0)],
            ["Columns", summary.get("columns", 0)],
            ["Numeric columns", ", ".join(summary.get("numeric_columns", []))],
            ["Datetime columns", ", ".join(summary.get("datetime_columns", []))],
        ]
        pd.DataFrame(summary_rows, columns=["Metric", "Value"]).to_excel(
            writer, sheet_name="Summary", index=False
        )
    return out.getvalue()
